"""Office document parsers and LibreOffice conversion helper."""

from __future__ import annotations

import asyncio
import os
import shutil
import signal
import subprocess
import tempfile
from collections.abc import AsyncIterator
from contextlib import asynccontextmanager, suppress
from pathlib import Path

from app.common.config import settings as app_settings
from app.common.exceptions import AppException
from app.common.schemas import ErrorCode
from app.parsers.base import DocumentParser, ParsedAsset, ParsedBlock, ParsedDocument

_PROCESS_STOP_TIMEOUT_SECONDS = 5.0


async def _start_libreoffice(*cmd: str) -> asyncio.subprocess.Process:
    if os.name == "posix":
        return await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            start_new_session=True,
        )
    if os.name == "nt":
        return await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            creationflags=getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0),
        )
    return await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )


async def _kill_libreoffice_tree(proc: asyncio.subprocess.Process) -> None:
    if os.name == "posix":
        with suppress(ProcessLookupError):
            os.killpg(proc.pid, signal.SIGKILL)
    elif os.name == "nt":
        try:
            killer = await asyncio.create_subprocess_exec(
                "taskkill",
                "/PID",
                str(proc.pid),
                "/T",
                "/F",
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL,
            )
            await asyncio.wait_for(
                killer.communicate(), timeout=_PROCESS_STOP_TIMEOUT_SECONDS
            )
        except (OSError, asyncio.TimeoutError):
            pass


async def _terminate_libreoffice(proc: asyncio.subprocess.Process) -> None:
    """Terminate the dedicated conversion process tree and reap the child."""
    if proc.returncode is not None:
        return
    await _kill_libreoffice_tree(proc)
    if proc.returncode is None:
        with suppress(ProcessLookupError):
            proc.kill()
    try:
        await asyncio.wait_for(
            proc.communicate(), timeout=_PROCESS_STOP_TIMEOUT_SECONDS
        )
    except asyncio.TimeoutError:
        if proc.returncode is None:
            with suppress(ProcessLookupError):
                proc.kill()
        with suppress(asyncio.TimeoutError):
            await asyncio.wait_for(proc.wait(), timeout=_PROCESS_STOP_TIMEOUT_SECONDS)


@asynccontextmanager
async def convert_with_libreoffice(
    source: Path, target_ext: str = "docx"
) -> AsyncIterator[Path]:
    """Convert in a private temporary directory owned by the caller context."""
    cfg = app_settings
    soffice = shutil.which(cfg.libreoffice_bin) or shutil.which("libreoffice")
    if not soffice:
        raise AppException(
            code=ErrorCode.CONVERSION_FAILED,
            message="LibreOffice 未安装，无法转换旧格式文档",
            status_code=500,
        )
    with tempfile.TemporaryDirectory(prefix="lo-convert-") as temp_dir:
        out_dir = Path(temp_dir)
        profile_dir = out_dir / "profile"
        profile_dir.mkdir()
        cmd = [
            soffice,
            f"-env:UserInstallation={profile_dir.as_uri()}",
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--norestore",
            "--nodefault",
            "--nofirststartwizard",
            "--convert-to",
            target_ext,
            "--outdir",
            str(out_dir),
            str(source),
        ]
        try:
            proc = await _start_libreoffice(*cmd)
        except OSError as exc:
            raise AppException(
                code=ErrorCode.CONVERSION_FAILED,
                message="LibreOffice 转换启动失败",
                status_code=500,
            ) from exc
        try:
            await asyncio.wait_for(
                proc.communicate(), timeout=cfg.libreoffice_timeout_seconds
            )
        except asyncio.TimeoutError as exc:
            await _terminate_libreoffice(proc)
            raise AppException(
                code=ErrorCode.CONVERSION_FAILED,
                message="LibreOffice 转换超时",
                status_code=500,
            ) from exc
        except BaseException:
            await _terminate_libreoffice(proc)
            raise
        if proc.returncode != 0:
            raise AppException(
                code=ErrorCode.CONVERSION_FAILED,
                message="LibreOffice 转换失败",
                status_code=500,
            )
        outputs = list(out_dir.glob(f"*.{target_ext}"))
        if not outputs:
            raise AppException(
                code=ErrorCode.CONVERSION_FAILED,
                message="LibreOffice 未产出目标文件",
                status_code=500,
            )
        yield outputs[0]


class DocxParser(DocumentParser):
    name = "docx"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".docx" or "wordprocessingml" in mime_type

    async def parse(self, source_path: str) -> ParsedDocument:
        from docx import Document as DocxDocument

        path = Path(source_path)
        doc = DocxDocument(str(path))
        blocks: list[ParsedBlock] = []
        assets: list[ParsedAsset] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if "heading" in style:
                level = 1
                for ch in style:
                    if ch.isdigit():
                        level = int(ch)
                        break
                blocks.append(ParsedBlock(text=text, block_type="heading", level=level, page_no=1))
            else:
                blocks.append(ParsedBlock(text=text, block_type="paragraph", page_no=1))
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
            if not rows:
                continue
            md = [
                "| " + " | ".join(rows[0]) + " |",
                "| " + " | ".join("---" for _ in rows[0]) + " |",
            ]
            for row in rows[1:]:
                md.append("| " + " | ".join(row) + " |")
            blocks.append(ParsedBlock(text="\n".join(md), block_type="table", page_no=1))
        # Extract images from relationships when available
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                blob = rel.target_part.blob
                name = Path(rel.target_ref).name
                mime = "image/png" if name.lower().endswith(".png") else "image/jpeg"
                assets.append(ParsedAsset(filename=name, data=blob, mime_type=mime, page_no=1))
        return ParsedDocument(
            title=path.stem,
            blocks=blocks or [ParsedBlock(text="", block_type="paragraph", page_no=1)],
            assets=assets,
            page_count=1,
            parser_name=self.name,
        )


class DocParser(DocumentParser):
    name = "doc"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".doc" or mime_type == "application/msword"

    async def parse(self, source_path: str) -> ParsedDocument:
        async with convert_with_libreoffice(Path(source_path), "docx") as converted:
            parsed = await DocxParser().parse(str(converted))
        parsed.parser_name = self.name
        parsed.warnings.append("converted_via_libreoffice")
        return parsed


class XlsxParser(DocumentParser):
    name = "xlsx"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".xlsx" or "spreadsheetml" in mime_type

    async def parse(self, source_path: str) -> ParsedDocument:
        from openpyxl import load_workbook

        path = Path(source_path)
        wb = load_workbook(str(path), data_only=True, read_only=True)
        blocks: list[ParsedBlock] = []
        page = 0
        for sheet in wb.worksheets:
            page += 1
            blocks.append(
                ParsedBlock(text=sheet.title, block_type="heading", level=2, page_no=page, sheet_name=sheet.title)
            )
            rows = []
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i > 500:
                    blocks.append(
                        ParsedBlock(
                            text=f"（工作表 {sheet.title} 超过 500 行，完整数据保留在原件中）",
                            block_type="paragraph",
                            page_no=page,
                            sheet_name=sheet.title,
                        )
                    )
                    break
                values = [("" if c is None else str(c)) for c in row]
                if any(v.strip() for v in values):
                    rows.append(values)
            if rows:
                width = max(len(r) for r in rows)
                rows = [r + [""] * (width - len(r)) for r in rows]
                md = [
                    "| " + " | ".join(rows[0]) + " |",
                    "| " + " | ".join("---" for _ in range(width)) + " |",
                ]
                for row in rows[1:]:
                    md.append("| " + " | ".join(row) + " |")
                blocks.append(
                    ParsedBlock(
                        text="\n".join(md),
                        block_type="table",
                        page_no=page,
                        sheet_name=sheet.title,
                    )
                )
        wb.close()
        return ParsedDocument(
            title=path.stem,
            blocks=blocks or [ParsedBlock(text="(空工作簿)", block_type="paragraph", page_no=1)],
            page_count=page or 1,
            parser_name=self.name,
        )


class XlsParser(DocumentParser):
    name = "xls"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".xls" or mime_type == "application/vnd.ms-excel"

    async def parse(self, source_path: str) -> ParsedDocument:
        async with convert_with_libreoffice(Path(source_path), "xlsx") as converted:
            parsed = await XlsxParser().parse(str(converted))
        parsed.parser_name = self.name
        parsed.warnings.append("converted_via_libreoffice")
        return parsed


class PptxParser(DocumentParser):
    name = "pptx"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".pptx" or "presentationml" in mime_type

    async def parse(self, source_path: str) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        path = Path(source_path)
        prs = Presentation(str(path))
        blocks: list[ParsedBlock] = []
        assets: list[ParsedAsset] = []
        for idx, slide in enumerate(prs.slides, start=1):
            blocks.append(
                ParsedBlock(
                    text=f"幻灯片 {idx}",
                    block_type="heading",
                    level=2,
                    page_no=idx,
                    slide_no=idx,
                )
            )
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        blocks.append(
                            ParsedBlock(text=text, block_type="paragraph", page_no=idx, slide_no=idx)
                        )
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    ext = image.ext or "png"
                    name = f"slide-{idx}-{len(assets)+1}.{ext}"
                    assets.append(
                        ParsedAsset(
                            filename=name,
                            data=image.blob,
                            mime_type=image.content_type or "image/png",
                            page_no=idx,
                        )
                    )
            # slide notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(
                        ParsedBlock(
                            text=f"备注: {notes}",
                            block_type="paragraph",
                            page_no=idx,
                            slide_no=idx,
                        )
                    )
        return ParsedDocument(
            title=path.stem,
            blocks=blocks or [ParsedBlock(text="", block_type="paragraph", page_no=1)],
            assets=assets,
            page_count=len(prs.slides),
            parser_name=self.name,
        )


class PptParser(DocumentParser):
    name = "ppt"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension == ".ppt" or mime_type == "application/vnd.ms-powerpoint"

    async def parse(self, source_path: str) -> ParsedDocument:
        async with convert_with_libreoffice(Path(source_path), "pptx") as converted:
            parsed = await PptxParser().parse(str(converted))
        parsed.parser_name = self.name
        parsed.warnings.append("converted_via_libreoffice")
        return parsed


class OpenDocumentParser(DocumentParser):
    name = "opendocument"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension in {".odt", ".ods", ".odp", ".rtf"} or "opendocument" in mime_type

    async def parse(self, source_path: str) -> ParsedDocument:
        path = Path(source_path)
        ext = path.suffix.lower()
        target = {"odt": "docx", "ods": "xlsx", "odp": "pptx", "rtf": "docx"}.get(ext.lstrip("."), "docx")
        async with convert_with_libreoffice(path, target) as converted:
            if target == "docx":
                parsed = await DocxParser().parse(str(converted))
            elif target == "xlsx":
                parsed = await XlsxParser().parse(str(converted))
            else:
                parsed = await PptxParser().parse(str(converted))
        parsed.parser_name = self.name
        parsed.warnings.append("converted_via_libreoffice")
        return parsed
