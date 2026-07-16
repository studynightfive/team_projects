#!/usr/bin/env python3
# ============================================================
# 智能知识库平台 - 二进制样本生成脚本
# 用途：生成 PDF、DOCX、XLSX 等二进制格式的测试样本
# 用法：uv run python scripts/generate_binary_samples.py
# 注意：需要安装 python-docx、openpyxl、fpdf 等库
# ============================================================

"""
本脚本用于生成二进制格式的测试样本文件。

使用方法：
  1. 安装依赖：pip install python-docx openpyxl fpdf2 Pillow
  2. 运行：python scripts/generate_binary_samples.py
  3. 生成的文件将放在 samples/documents/ 对应目录下

注意：生成的文件内容为假数据，不包含真实个人信息。
"""

import os
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SAMPLES_DIR = PROJECT_ROOT / "samples" / "documents"


def generate_docx_samples():
    """生成 Word 文档样本"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
    except ImportError:
        print("[跳过] python-docx 未安装，无法生成 Word 样本")
        return

    # 1. 简单文本 DOCX
    doc = Document()
    doc.styles["Normal"].font.size = Pt(12)
    doc.add_heading("智能知识库平台", level=1)
    doc.add_paragraph("版本：0.1.0")
    doc.add_paragraph("日期：2026年7月15日")
    doc.add_heading("项目概述", level=2)
    doc.add_paragraph(
        "智能知识库平台是一个面向团队内部文档集中管理、检索、问答和引用核验的综合平台。"
        "平台采用前后端分离架构，前端使用Vue 3和TypeScript，后端基于FastAPI提供高性能的API服务。"
    )
    doc.add_heading("核心功能", level=2)
    for feature in [
        "多格式文档上传与自动转换",
        "OCR文字识别与版面分析",
        "关键词、向量和混合检索",
        "流式问答与引用溯源",
        "文档五格式导出",
    ]:
        doc.add_paragraph(feature, style="List Bullet")
    doc.save(str(SAMPLES_DIR / "word" / "normal" / "simple.docx"))
    print("[生成] word/normal/simple.docx")

    # 2. 含表格 DOCX
    doc = Document()
    doc.add_heading("季度销售数据", level=1)
    table = doc.add_table(rows=4, cols=5, style="Table Grid")
    headers = ["产品名称", "一月", "二月", "三月", "合计"]
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header
    data = [
        ["知识库平台", "120", "135", "150", "405"],
        ["文档处理引擎", "85", "92", "88", "265"],
        ["检索增强模块", "65", "70", "78", "213"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx + 1].cells[col_idx].text = cell_data
    doc.save(str(SAMPLES_DIR / "word" / "normal" / "with-table.docx"))
    print("[生成] word/normal/with-table.docx")

    # 3. 含图片 DOCX
    doc = Document()
    doc.add_heading("含图片的文档", level=1)
    doc.add_paragraph("以下是一张示例图片：")
    # 创建一个简单的占位图片
    img_path = PROJECT_ROOT / "temp_sample_image.png"
    try:
        from PIL import Image
        img = Image.new("RGB", (200, 100), color=(73, 109, 137))
        img.save(img_path)
        doc.add_picture(str(img_path), width=Inches(3))
        img_path.unlink()  # 删除临时图片
    except ImportError:
        doc.add_paragraph("[图片占位 - Pillow未安装]")
    doc.save(str(SAMPLES_DIR / "word" / "normal" / "with-images.docx"))
    print("[生成] word/normal/with-images.docx")


def generate_xlsx_samples():
    """生成 Excel 文档样本"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment
    except ImportError:
        print("[跳过] openpyxl 未安装，无法生成 Excel 样本")
        return

    # 1. 简单 XLSX
    wb = Workbook()
    ws = wb.active
    ws.title = "数据表"
    ws["A1"] = "产品名称"
    ws["B1"] = "数量"
    ws["C1"] = "单价"
    ws["D1"] = "总价"
    data = [
        ["知识库平台", 120, 5000, 600000],
        ["文档处理引擎", 85, 3000, 255000],
        ["检索增强模块", 65, 4000, 260000],
    ]
    for row in data:
        ws.append(row)
    wb.save(str(SAMPLES_DIR / "excel" / "normal" / "simple.xlsx"))
    print("[生成] excel/normal/simple.xlsx")

    # 2. 合并单元格 XLSX
    wb = Workbook()
    ws = wb.active
    ws.title = "合并单元格"
    ws.merge_cells("A1:D1")
    ws["A1"] = "2026年第一季度销售数据"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A1"].alignment = Alignment(horizontal="center")
    ws["A3"] = "产品"
    ws["B3"] = "一月"
    ws["C3"] = "二月"
    ws["D3"] = "三月"
    for i, row in enumerate(data, start=4):
        ws.append(row[:4])
    ws.merge_cells("A6:D6")
    ws["A6"] = "合计：1,115,000元"
    wb.save(str(SAMPLES_DIR / "excel" / "normal" / "merged-cells.xlsx"))
    print("[生成] excel/normal/merged-cells.xlsx")

    # 3. 多工作表 XLSX
    wb = Workbook()
    for sheet_name in ["Q1", "Q2", "Q3"]:
        ws = wb.create_sheet(sheet_name)
        ws["A1"] = f"{sheet_name} 数据"
        ws.append(["项目", "数值"])
        ws.append(["产品A", 100])
        ws.append(["产品B", 200])
    wb.remove(wb["Sheet"])  # 删除默认工作表
    wb.save(str(SAMPLES_DIR / "excel" / "normal" / "multi-sheet.xlsx"))
    print("[生成] excel/normal/multi-sheet.xlsx")

    # 4. 空工作簿 XLSX
    wb = Workbook()
    wb.save(str(SAMPLES_DIR / "excel" / "abnormal" / "empty.xlsx"))
    print("[生成] excel/abnormal/empty.xlsx")


def generate_pptx_samples():
    """生成 PowerPoint 样本"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("[跳过] python-pptx 未安装，无法生成 PPT 样本")
        return

    # 1. 简单 PPTX
    prs = Presentation()
    # 幻灯片 1：标题
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "智能知识库平台"
    slide.shapes.placeholders[1].text = "产品介绍\n2026年7月"
    # 幻灯片 2：内容
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "核心功能"
    content = slide.shapes.placeholders[1].text_frame
    content.text = "1. 多格式文档上传与转换\n2. OCR文字识别\n3. 智能检索\n4. 流式问答\n5. 文档导出"
    prs.save(str(SAMPLES_DIR / "powerpoint" / "normal" / "simple.pptx"))
    print("[生成] powerpoint/normal/simple.pptx")

    # 2. 含备注 PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "含备注的演示文稿"
    notes = slide.notes_slide
    notes.notes_text_frame.text = "这是演讲者备注内容，用于测试备注提取功能。"
    prs.save(str(SAMPLES_DIR / "powerpoint" / "normal" / "with-notes.pptx"))
    print("[生成] powerpoint/normal/with-notes.pptx")

    # 3. 无内容 PPTX
    prs = Presentation()
    prs.save(str(SAMPLES_DIR / "powerpoint" / "abnormal" / "no-content.pptx"))
    print("[生成] powerpoint/abnormal/no-content.pptx")


def generate_pdf_samples():
    """生成 PDF 样本"""
    try:
        from fpdf import FPDF
    except ImportError:
        print("[跳过] fpdf2 未安装，无法生成 PDF 样本")
        return

    # 1. 简单文本 PDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=16)
    pdf.cell(200, 10, text="Knowledge Base Platform", align="C")
    pdf.ln(15)
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, text="Version 0.1.0", align="C")
    pdf.ln(15)
    pdf.set_font("Helvetica", size=10)
    text = (
        "This is a simple PDF document for testing document parsing. "
        "The platform supports multiple document formats including PDF, "
        "Word, Excel, PowerPoint, Markdown, and more."
    )
    pdf.multi_cell(0, 10, text)
    pdf.output(str(SAMPLES_DIR / "pdf" / "normal" / "simple-text.pdf"))
    print("[生成] pdf/normal/simple-text.pdf")

    # 2. 含表格 PDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=16)
    pdf.cell(200, 10, text="Quarterly Report", align="C")
    pdf.ln(15)
    pdf.set_font("Helvetica", size=10)
    headers = ["Product", "Jan", "Feb", "Mar", "Total"]
    col_width = 40
    for header in headers:
        pdf.cell(col_width, 10, header, border=1)
    pdf.ln()
    data = [
        ["Platform", "120", "135", "150", "405"],
        ["Engine", "85", "92", "88", "265"],
        ["Module", "65", "70", "78", "213"],
    ]
    for row in data:
        for cell in row:
            pdf.cell(col_width, 10, cell, border=1)
        pdf.ln()
    pdf.output(str(SAMPLES_DIR / "pdf" / "normal" / "with-table.pdf"))
    print("[生成] pdf/normal/with-table.pdf")

    # 3. 损坏 PDF（截断）
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    pdf.cell(200, 10, text="This PDF is intentionally corrupted.")
    pdf.output(str(SAMPLES_DIR / "pdf" / "abnormal" / "corrupted.pdf"))
    # 截断文件
    with open(SAMPLES_DIR / "pdf" / "abnormal" / "corrupted.pdf", "rb") as f:
        data = f.read()
    with open(SAMPLES_DIR / "pdf" / "abnormal" / "corrupted.pdf", "wb") as f:
        f.write(data[:len(data) // 2])  # 只保留前半部分
    print("[生成] pdf/abnormal/corrupted.pdf")


def main():
    """主流程"""
    print("=" * 60)
    print("  智能知识库平台 - 二进制样本生成")
    print("=" * 60)
    print()

    generate_docx_samples()
    print()
    generate_xlsx_samples()
    print()
    generate_pptx_samples()
    print()
    generate_pdf_samples()
    print()

    print("=" * 60)
    print("  生成完成！")
    print("  请运行 python scripts/validate_samples.py 验证")
    print("=" * 60)


if __name__ == "__main__":
    main()